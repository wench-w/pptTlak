from __future__ import annotations

import re
import subprocess
from pathlib import Path
from typing import Iterable, List

import fitz  # type: ignore
from pptx import Presentation  # type: ignore

from .models import SlideAsset

PNG_PATTERN = re.compile(r"(?:\D|^)(\d+)(?=\.png$)", re.IGNORECASE)


def convert_ppt_to_images(
    ppt_path: Path,
    output_dir: Path,
    libreoffice: str = "libreoffice",
) -> List[Path]:
    """
    调用 LibreOffice 将 PPT/PPTX 导出为 PNG。

    Returns:
        生成的 PNG 路径列表（按页序排序）。
    """

    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_output = output_dir / f"{ppt_path.stem}.pdf"
    cmd = [
        libreoffice,
        "--headless",
        "--convert-to",
        "pdf:impress_pdf_Export",
        "--outdir",
        str(output_dir),
        str(ppt_path),
    ]
    subprocess.run(cmd, check=True)
    if not pdf_output.exists():
        raise RuntimeError(f"LibreOffice 未生成 PDF: {pdf_output}")
    render_pdf_to_png(pdf_output, output_dir, ppt_path.stem)
    pdf_output.unlink(missing_ok=True)
    return sorted_pngs(output_dir)


def render_pdf_to_png(pdf_path: Path, output_dir: Path, prefix: str) -> None:
    """将PDF渲染为PNG，使用高DPI确保字体清晰。"""
    doc = fitz.open(pdf_path)
    for idx in range(doc.page_count):
        page = doc[idx]
        # 提高DPI到300，确保PPT内容字体清晰
        pix = page.get_pixmap(dpi=300)
        out_path = output_dir / f"{prefix}_{idx + 1:03d}.png"
        pix.save(out_path)
    doc.close()


def sorted_pngs(directory: Path) -> List[Path]:
    def sort_key(path: Path) -> int:
        match = PNG_PATTERN.search(path.name)
        return int(match.group(1)) if match else 0

    return sorted(directory.glob("*.png"), key=sort_key)


def extract_slide_notes(ppt_path: Path) -> tuple[List[str], List[tuple[str | None, int | None]]]:
    """
    提取每页幻灯片的备注文本和字体信息，支持多段落。
    
    Returns:
        (notes, font_info): notes是文本列表，font_info是(字体名, 字号)列表
    """
    prs = Presentation(str(ppt_path))
    notes = []
    font_info_list = []
    
    for slide in prs.slides:
        text = ""
        font_name = None
        font_size = None
        
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            # 获取备注文本框架中的所有文本
            notes_frame = slide.notes_slide.notes_text_frame
            # 如果有多段落，遍历所有段落
            if notes_frame.paragraphs:
                text_parts = []
                for paragraph in notes_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_parts.append(para_text)
                        # 尝试从段落中提取字体信息
                        if paragraph.runs:
                            # 使用第一个run的字体信息（通常备注使用统一字体）
                            first_run = paragraph.runs[0]
                            if first_run.font and first_run.font.name:
                                font_name = first_run.font.name
                            if first_run.font and first_run.font.size:
                                # font.size是EMU单位，转换为磅
                                font_size = int(first_run.font.size / 12700) if first_run.font.size else None
                text = "\n".join(text_parts)
            else:
                # 如果没有段落结构，直接获取文本
                text = notes_frame.text
                # 尝试从文本框架的默认格式获取字体信息
                if hasattr(notes_frame, 'paragraphs') and notes_frame.paragraphs:
                    first_para = notes_frame.paragraphs[0]
                    if first_para.runs:
                        first_run = first_para.runs[0]
                        if first_run.font and first_run.font.name:
                            font_name = first_run.font.name
                        if first_run.font and first_run.font.size:
                            font_size = int(first_run.font.size / 12700) if first_run.font.size else None
        
        notes.append(text.strip() if text else "")
        font_info_list.append((font_name, font_size))
    
    return notes, font_info_list


def build_slide_assets(
    images: Iterable[Path],
    notes: Iterable[str],
    font_info: Iterable[tuple[str | None, int | None]] | None = None,
) -> List[SlideAsset]:
    """构建幻灯片资源，包括备注文本和字体信息。"""
    assets: List[SlideAsset] = []
    font_iter = iter(font_info) if font_info else iter([(None, None)] * len(list(images)))
    
    for idx, (image_path, note) in enumerate(zip(images, notes), start=1):
        font_name, font_size = next(font_iter, (None, None))
        assets.append(
            SlideAsset(
                index=idx,
                image_path=image_path,
                note_text=note or "",
                note_font=font_name,
                note_fontsize=font_size,
            )
        )
    return assets

